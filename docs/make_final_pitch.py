# -*- coding: utf-8 -*-
"""Generate the 10-minute championship pitch deck for EveryLane Macau."""
from __future__ import annotations

import os
from pathlib import Path

import qrcode
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
SEMI = ASSETS / "semifinal"
OUT = HERE / "決賽路演_10分鐘_街知巷聞.pptx"
QR = SEMI / "site_qr.png"
SITE = "http://47.79.228.128/"

CREAM = RGBColor(249, 244, 233)
PAPER = RGBColor(255, 253, 248)
INK = RGBColor(42, 32, 25)
MUTED = RGBColor(105, 93, 80)
TERRA = RGBColor(190, 74, 58)
NAVY = RGBColor(38, 63, 82)
BLUE = RGBColor(44, 94, 134)
GOLD = RGBColor(194, 145, 46)
GREEN = RGBColor(62, 142, 90)
LINE = RGBColor(225, 213, 193)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def fill_bg(slide, color=CREAM):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def shape(slide, x, y, w, h, fill, line=None, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line or fill
    return s


def textbox(slide, text, x, y, w, h, size=20, color=INK, bold=False,
            font="Noto Sans TC", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
            margin=0.02):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return b


def rich_lines(slide, lines, x, y, w, h, size=18, color=INK, gap=8):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = b.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            head, body = item
            r = p.add_run()
            r.text = head
            r.font.bold = True
            r.font.color.rgb = TERRA
            r2 = p.add_run()
            r2.text = body
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.color.rgb = color
        for run in p.runs:
            run.font.name = "Noto Sans TC"
            run.font.size = Pt(size)
        p.space_after = Pt(gap)
        p.line_spacing = 1.14
    return b


def title(slide, kicker, heading, number):
    textbox(slide, kicker, .65, .38, 9.8, .32, 10, BLUE, True)
    textbox(slide, heading, .65, .78, 11.7, .72, 28, INK, True, "Noto Serif TC")
    textbox(slide, f"{number:02d}", 12.25, .42, .45, .35, 10, MUTED, True,
            align=PP_ALIGN.RIGHT)
    shape(slide, .65, 1.55, 1.05, .04, TERRA, radius=False)


def footer(slide):
    textbox(slide, "街知巷聞 EveryLane Macau  ·  愛拼才會贏  ·  施天益",
            .65, 7.12, 8.5, .22, 8, MUTED)
    textbox(slide, SITE, 9.6, 7.12, 3.05, .22, 8, MUTED,
            align=PP_ALIGN.RIGHT)


def add_image(slide, path, x, y, w, h=None):
    path = str(path)
    if not os.path.exists(path):
        return None
    if h is None:
        with Image.open(path) as im:
            h = w * im.height / im.width
    return slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w),
                                    height=Inches(h))


def metric(slide, value, label, x, y, w=1.75, color=TERRA):
    shape(slide, x, y, w, 1.08, PAPER, LINE)
    textbox(slide, value, x + .12, y + .15, w - .24, .4, 23, color, True,
            "Noto Serif TC")
    textbox(slide, label, x + .12, y + .66, w - .24, .22, 9, MUTED, True)


def pill(slide, text, x, y, w, fill, color):
    shape(slide, x, y, w, .38, fill, fill)
    textbox(slide, text, x + .08, y + .08, w - .16, .18, 9, color, True,
            align=PP_ALIGN.CENTER)


def new_slide(bg=CREAM):
    s = prs.slides.add_slide(blank)
    fill_bg(s, bg)
    return s


# 1 — cover
s = new_slide()
shape(s, 0, 0, 13.333, 7.5, NAVY, radius=False)
textbox(s, "「千模百煉」AI 開發者系列學生競賽", .72, .55, 7.8, .35,
        11, RGBColor(224, 205, 159), True)
textbox(s, "街知巷聞", .72, 1.35, 5.2, .75, 38, WHITE, True, "Noto Serif TC")
textbox(s, "EveryLane Macau", .74, 2.14, 5.4, .38, 16,
        RGBColor(208, 220, 228), True)
textbox(s, "不只推薦澳門，\n而是重新分配澳門的旅遊流量。", .72, 2.85, 7.2, 1.25,
        28, RGBColor(243, 209, 132), True, "Noto Serif TC")
textbox(s, "Qwen / QwenPaw 智能體 × 舊區導流 × 可歸因商戶閉環",
        .74, 4.42, 7.4, .36, 14, WHITE, True)
textbox(s, "隊伍：愛拼才會贏  ·  施天益（SI TIN IEK）", .74, 5.12,
        6.8, .3, 12, RGBColor(208, 220, 228))
qrcode.make(SITE).save(QR)
shape(s, 10.05, 1.35, 2.45, 3.7, PAPER, PAPER)
add_image(s, QR, 10.4, 1.72, 1.75, 1.75)
textbox(s, "掃碼即場體驗", 10.25, 3.68, 2.05, .3, 13, NAVY, True,
        align=PP_ALIGN.CENTER)
textbox(s, "真實 Qwen\n90 秒評審模式\n成效儀表板", 10.3, 4.08, 1.95, .78,
        11, MUTED, False, align=PP_ALIGN.CENTER)
textbox(s, SITE, 9.55, 6.75, 3.0, .25, 10, WHITE, False,
        align=PP_ALIGN.RIGHT)

# 2 — problem
s = new_slide()
title(s, "01 · REAL PROBLEM", "澳門不缺景點，缺的是「把人流帶對地方」", 2)
shape(s, .68, 1.9, 3.6, 3.95, RGBColor(252, 237, 230), RGBColor(235, 192, 181))
textbox(s, "熱門點", .95, 2.2, 2.9, .35, 20, TERRA, True, "Noto Serif TC")
textbox(s, "大三巴／議事亭／路氹\n擠迫、排隊、同質體驗", .95, 2.82,
        2.9, 1.0, 17, INK, True)
textbox(s, "遊客：體驗下降\n城市：承載失衡", .95, 4.22, 2.9, .75,
        14, MUTED)
textbox(s, "→", 4.55, 3.35, .65, .55, 30, GOLD, True,
        align=PP_ALIGN.CENTER)
shape(s, 5.45, 1.9, 3.6, 3.95, RGBColor(237, 246, 241), RGBColor(181, 218, 195))
textbox(s, "舊區與小店", 5.72, 2.2, 2.9, .35, 20, GREEN, True, "Noto Serif TC")
textbox(s, "福隆新街／內港／下環\n有文化、有故事、缺客流", 5.72, 2.82,
        2.95, 1.0, 17, INK, True)
textbox(s, "社區：活力流失\n商戶：難量度導流價值", 5.72, 4.22, 2.95, .75,
        14, MUTED)
shape(s, 9.45, 1.9, 3.15, 3.95, NAVY, NAVY)
textbox(s, "AI 的任務", 9.75, 2.2, 2.55, .35, 20,
        RGBColor(243, 209, 132), True, "Noto Serif TC")
rich_lines(s, [
    ("不是：", "再做一個熱門榜單"),
    ("而是：", "在限制條件下主動導流"),
    ("最後：", "量度到訪與轉化"),
], 9.75, 2.9, 2.45, 1.8, 15, WHITE, 12)
textbox(s, "遊客 × 小店 × 城市三方共贏", .68, 6.32, 11.9, .38,
        19, TERRA, True, "Noto Serif TC", PP_ALIGN.CENTER)
footer(s)

# 3 — product
s = new_slide()
title(s, "02 · PRODUCT", "一句需求 → 一份可驗證、可執行、可導流的行程", 3)
add_image(s, SEMI / "05_itinerary_access_code.png", .65, 1.82, 7.55, 4.72)
rich_lines(s, [
    ("理解：", "日期、人數、興趣、預算、步行偏好"),
    ("執行：", "查天氣、開放、人流、路線與預算"),
    ("恢復：", "休息日換點；擁擠時導向附近老街"),
    ("驗證：", "地圖、時間軸、條件核對、無障礙資訊"),
    ("歸因：", "本地商戶站點可領一次性到店碼"),
], 8.55, 1.95, 4.0, 3.35, 16, INK, 10)
pill(s, "90 秒評審快速演示", 8.6, 5.62, 2.25,
     RGBColor(245, 231, 200), RGBColor(92, 65, 25))
pill(s, "自訂問題：真實 Qwen", 10.98, 5.62, 1.55,
     RGBColor(230, 241, 248), BLUE)
footer(s)

# 4 — architecture
s = new_slide()
title(s, "03 · AGENT DEPTH", "Qwen 做決策，工具做事實，安全層決定能否提交", 4)
layers = [
    ("Qwen / QwenPaw", "理解需求 · ReAct 決策 · function calling", NAVY, WHITE),
    ("EveryLane Skill", "阿濠人設 · 舊區導流目標 · 失敗恢復規則", BLUE, WHITE),
    ("stdio MCP 7+1", "搜尋 · 天氣 · 開放 · 人流 · 本地點 · 路線 · 預算", TERRA, WHITE),
    ("確定性核對器", "同區 · 當日開放 · 少步行 · 預算 · 結構化提交", GOLD, INK),
    ("產品輸出", "SSE 軌跡 · 地圖時間軸 · 到店碼 · B 端 JSON", GREEN, WHITE),
]
for i, (head, body, color, txt) in enumerate(layers):
    y = 1.84 + i * .88
    shape(s, 1.0, y, 11.3, .66, color, color)
    textbox(s, head, 1.25, y + .14, 2.25, .25, 15, txt, True)
    textbox(s, body, 3.65, y + .14, 8.25, .25, 13, txt)
    if i < len(layers) - 1:
        textbox(s, "↓", 6.2, y + .66, .8, .2, 13, MUTED, True,
                align=PP_ALIGN.CENTER)
textbox(s, "核心原則：模型不能杜撰開放、人流、距離與預算；未完成工具核對就不能 submit。",
        1.0, 6.47, 11.3, .35, 15, TERRA, True, align=PP_ALIGN.CENTER)
footer(s)

# 5 — recovery
s = new_slide()
title(s, "04 · FAILURE RECOVERY", "智能體的價值，不在順境推薦，而在限制衝突時仍完成任務", 5)
shape(s, .75, 1.88, 3.3, 3.8, PAPER, LINE)
textbox(s, "使用者要求", 1.0, 2.15, 2.8, .3, 16, NAVY, True)
textbox(s, "「星期三想去鄭家大屋\n同附近歷史老街」", 1.0, 2.75,
        2.75, .85, 20, INK, True, "Noto Serif TC")
pill(s, "日期 + 指名 POI + 興趣", 1.0, 4.42, 2.2,
     RGBColor(231, 240, 246), BLUE)
textbox(s, "→", 4.18, 3.25, .55, .55, 28, GOLD, True,
        align=PP_ALIGN.CENTER)
shape(s, 4.85, 1.88, 3.3, 3.8, RGBColor(252, 237, 230), RGBColor(235, 192, 181))
textbox(s, "工具發現衝突", 5.1, 2.15, 2.8, .3, 16, TERRA, True)
textbox(s, "check_opening\n鄭家大屋：逢週三休息", 5.1, 2.82,
        2.7, .85, 19, INK, True, "Noto Serif TC")
pill(s, "不可直接提交", 5.1, 4.42, 1.7,
     RGBColor(248, 220, 215), TERRA)
textbox(s, "→", 8.28, 3.25, .55, .55, 28, GOLD, True,
        align=PP_ALIGN.CENTER)
shape(s, 8.95, 1.88, 3.62, 3.8, RGBColor(237, 246, 241), RGBColor(181, 218, 195))
textbox(s, "自動恢復並重算", 9.2, 2.15, 3.0, .3, 16, GREEN, True)
textbox(s, "同區找開放替代點\n重新算路線、人流與預算\n最後才提交", 9.2, 2.75,
        2.95, 1.25, 18, INK, True, "Noto Serif TC")
pill(s, "任務仍完整完成", 9.2, 4.58, 2.1,
     RGBColor(218, 239, 226), GREEN)
textbox(s, "規劃 → 工具 → 觀察 → 修正 → 提交", .75, 6.2, 11.8, .38,
        20, TERRA, True, "Noto Serif TC", PP_ALIGN.CENTER)
footer(s)

# 6 — attribution loop
s = new_slide()
title(s, "05 · FROM RECOMMENDATION TO ATTRIBUTION", "第一次把「推薦舊區」做成可核銷的商業閉環", 6)
add_image(s, SEMI / "04_dashboard_redeem.png", .65, 1.82, 7.65, 4.75)
steps = [
    ("1", "行程導流", "本地商戶成為行程站點"),
    ("2", "遊客領碼", "EL-XXXX-XX 一次性碼"),
    ("3", "到店核銷", "第一次成功，第二次拒絕"),
    ("4", "效果結算", "由曝光轉為可量度到訪"),
]
for i, (n, head, body) in enumerate(steps):
    y = 1.92 + i * 1.08
    shape(s, 8.65, y, .48, .48, TERRA, TERRA)
    textbox(s, n, 8.65, y + .1, .48, .18, 11, WHITE, True,
            align=PP_ALIGN.CENTER)
    textbox(s, head, 9.35, y, 2.7, .25, 15, INK, True)
    textbox(s, body, 9.35, y + .35, 2.85, .3, 11, MUTED)
textbox(s, "功能是真實可操作；累計轉化數據按賽規使用示範數據並明確標註。",
        8.65, 6.2, 3.7, .48, 11, TERRA, True)
footer(s)

# 7 — metrics
s = new_slide()
title(s, "06 · IMPACT", "計劃書 9/9 項逐一呈現，但不把示範數據冒充田野研究", 7)
pill(s, "以下為複賽示範數據 · 口徑、公式、種子可下載核驗", .72, 1.76, 4.5,
     RGBColor(255, 243, 219), RGBColor(116, 83, 24))
metric(s, "86.4%", "導流覆蓋率", .72, 2.42, 2.15)
metric(s, "98.9%", "路線可行率", 3.02, 2.42, 2.15)
metric(s, "41.8%", "商戶到訪率", 5.32, 2.42, 2.15)
metric(s, "91.3%", "任務完成率", 7.62, 2.42, 2.15)
metric(s, "4.2", "平均舊區/商戶點", 9.92, 2.42, 2.15)
shape(s, .72, 4.05, 5.82, 1.55, RGBColor(237, 246, 241), RGBColor(181, 218, 195))
textbox(s, "真實・即時可核驗", .98, 4.32, 2.6, .3, 17, GREEN, True)
textbox(s, "Qwen / 工具鏈 · Open-Meteo · 到店碼 · B 端 API\n70 POI · 無障礙 · 自動化測試 · 系統 uptime",
        .98, 4.82, 5.0, .55, 13, INK)
shape(s, 6.78, 4.05, 5.55, 1.55, RGBColor(255, 243, 219), RGBColor(232, 210, 160))
textbox(s, "示範・確定性生成", 7.05, 4.32, 2.9, .3, 17,
        RGBColor(116, 83, 24), True)
textbox(s, "可用性 · 累計轉化 · 區域熱度 · 模型校正\n不宣稱真實；/api/impact/evidence 可審計",
        7.05, 4.82, 4.85, .55, 13, INK)
textbox(s, "透明不是減分項；把限制講清楚，才是可落地 AI 的基本功。",
        .72, 6.18, 11.6, .4, 18, NAVY, True, "Noto Serif TC",
        PP_ALIGN.CENTER)
footer(s)

# 8 — business
s = new_slide()
title(s, "07 · BUSINESS & MOAT", "一套導流能力，三條清晰變現路徑", 8)
cards = [
    ("商戶", "精選訂閱／按核銷效果付費", "到店碼是結算依據", TERRA),
    ("酒店／旅行社", "白標行程 API 授權", "已上線 v1 API + 歸因 JSON", BLUE),
    ("文旅部門", "匿名化區域熱度儀表板", "監測承載與舊區活化", GREEN),
]
for i, (head, model, proof, color) in enumerate(cards):
    x = .75 + i * 4.15
    shape(s, x, 1.95, 3.75, 2.65, PAPER, LINE)
    shape(s, x, 1.95, 3.75, .12, color, color, radius=False)
    textbox(s, head, x + .28, 2.28, 3.1, .34, 20, color, True,
            "Noto Serif TC")
    textbox(s, model, x + .28, 2.95, 3.1, .5, 16, INK, True)
    textbox(s, proof, x + .28, 3.82, 3.1, .35, 12, MUTED)
textbox(s, "為什麼不是一般旅遊聊天機械人？", .75, 5.1, 4.2, .35,
        19, NAVY, True, "Noto Serif TC")
rich_lines(s, [
    ("差異化資料：", "舊區/商戶、人流、開放、無障礙與故事層"),
    ("差異化決策：", "以承載平衡而非熱門度排序"),
    ("差異化閉環：", "推薦 → 到訪 → 核銷 → 成效歸因"),
], .75, 5.58, 11.5, .95, 14, INK, 4)
footer(s)

# 9 — engineering
s = new_slide()
title(s, "08 · TRUST & ENGINEERING", "路演不能靠運氣：真實 Qwen，也必須有可解釋的保護網", 9)
metric(s, "641", "後端 PASS", .75, 1.95, 2.05, BLUE)
metric(s, "119", "瀏覽器 PASS", 2.98, 1.95, 2.05, BLUE)
metric(s, "101", "API / 安全 PASS", 5.21, 1.95, 2.05, BLUE)
metric(s, "8", "MCP 工具協議 PASS", 7.44, 1.95, 2.05, BLUE)
metric(s, "0", "上述套件 FAIL", 9.67, 1.95, 2.05, GREEN)
shape(s, .75, 3.58, 5.7, 2.2, PAPER, LINE)
textbox(s, "可靠性", 1.02, 3.88, 2.0, .3, 18, TERRA, True)
rich_lines(s, [
    ("45 秒：", "單次模型逾時"),
    ("110 秒：", "整體 Qwen 時間預算"),
    ("自動接管：", "同一知識庫與工具完成，不顯示假 Qwen"),
    ("靜態備援：", "GitHub Pages 仍可演示核心流程"),
], 1.02, 4.35, 4.95, 1.15, 13, INK, 4)
shape(s, 6.75, 3.58, 5.58, 2.2, PAPER, LINE)
textbox(s, "AI 倫理", 7.02, 3.88, 2.0, .3, 18, GREEN, True)
rich_lines(s, [
    ("透明：", "引擎、估算、回退與示範數據明示"),
    ("準確：", "關鍵事實由工具與安全核對器控制"),
    ("私隱：", "免登入；到店碼不含個人資料"),
    ("公平：", "刻意讓資源較弱的舊區獲得曝光"),
], 7.02, 4.35, 4.85, 1.15, 13, INK, 4)
footer(s)

# 10 — close
s = new_slide()
shape(s, 0, 0, 13.333, 7.5, NAVY, radius=False)
textbox(s, "我們完成的，不是一份漂亮行程。", .8, .9, 11.7, .55,
        26, WHITE, True, "Noto Serif TC", PP_ALIGN.CENTER)
textbox(s, "而是一個能把旅遊流量帶入舊區、\n再用到店碼證明價值的 AI Agent。",
        1.25, 1.78, 10.8, 1.35, 32, RGBColor(243, 209, 132), True,
        "Noto Serif TC", PP_ALIGN.CENTER)
pill(s, "創意：舊區導流", 2.0, 3.65, 2.4, TERRA, WHITE)
pill(s, "技術：QwenPaw + MCP", 5.05, 3.65, 3.15, BLUE, WHITE)
pill(s, "價值：到店可歸因", 8.85, 3.65, 2.4, GREEN, WHITE)
textbox(s, "現場可驗證：90 秒評審演示 → 到店碼核銷 → 成效儀表板",
        1.05, 4.52, 11.2, .4, 17, WHITE, True, align=PP_ALIGN.CENTER)
add_image(s, QR, 5.75, 5.1, 1.45, 1.45)
textbox(s, SITE, 4.25, 6.72, 4.45, .26, 12, WHITE, True,
        align=PP_ALIGN.CENTER)
textbox(s, "多謝各位評審，歡迎提問。", 3.45, 7.08, 6.45, .26,
        11, RGBColor(208, 220, 228), align=PP_ALIGN.CENTER)

prs.save(OUT)
print("SAVED", OUT)
